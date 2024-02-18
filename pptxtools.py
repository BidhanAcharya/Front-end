from pptx import Presentation
from pptx.util import Inches, Pt
import re
def add_slide(prs, layout):
    return prs.slides.add_slide(layout)

def customizer_placeholder(slide,placeholder_number:int,font_name,font_size:int,bold:bool):
  slide.placeholders[placeholder_number].text_frame.paragraphs[0].font.name = font_name
  slide.placeholders[placeholder_number].text_frame.paragraphs[0].font.size = Pt(font_size)
  slide.placeholders[placeholder_number].text_frame.paragraphs[0].font.bold = bold

def customizer_placeholder_paragraphs(slide,placeholder_number:int,paragraph_number:int,font_name,font_size:int,bold:bool):
  slide.shapes.placeholders[placeholder_number].text_frame.paragraphs[paragraph_number].font.name = font_name
  slide.shapes.placeholders[placeholder_number].text_frame.paragraphs[paragraph_number].font.size = Pt(font_size)
  slide.shapes.placeholders[placeholder_number].text_frame.paragraphs[paragraph_number].font.bold = bold

def split_sentences(text):
  # Use regular expression to split sentences, excluding decimal digits
  sentences = re.split(r'(?<!\d\.\d.)(?<!\d\.\d)(?<![A-Z][a-z]\.)(?<!\w\.\w.)(?<=\.|\?)\s', text)
  return sentences

def shape_set_font_size(shape, points):
  for paragraph in shape.text_frame.paragraphs:
      paragraph.font.size = Pt(points)
